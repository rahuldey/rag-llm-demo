import logging
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

from rag_demo.domain.models import DocumentChunk
from rag_demo.domain.ports import DocumentLoaderPort

logger = logging.getLogger(__name__)

class GenericDocumentLoader(DocumentLoaderPort):
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200) -> None:
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )

    def load(self, file_path: Path) -> list[DocumentChunk]:
        suffix = file_path.suffix.lower()
        full_text = ""

        try:
            if suffix == ".pdf":
                reader = PdfReader(file_path)
                for page in reader.pages:
                    page_text = page.extract_text() or ""
                    full_text += f"\n{page_text}"
            elif suffix == ".txt":
                try:
                    with open(file_path, "r", encoding="utf-8") as f:
                        full_text = f.read()
                except UnicodeDecodeError:
                    with open(file_path, "r", encoding="latin-1") as f:
                        full_text = f.read()
            elif suffix in {".ppt", ".pptx"}:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            full_text += f"\n{shape.text}"
            else:
                logger.warning("Unsupported file type: %s", suffix)
                return []
        except Exception as e:
            logger.warning("Error reading %s: %s", file_path.name, e)
            return []

        if not full_text.strip():
            return []

        raw_chunks = self._splitter.split_text(full_text)

        return [
            DocumentChunk(
                content=chunk,
                metadata={
                    "source": file_path.name,
                    "chunk_index": str(i),
                },
            )
            for i, chunk in enumerate(raw_chunks)
        ]
